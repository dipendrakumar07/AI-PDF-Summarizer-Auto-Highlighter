import os
import fitz  
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import sent_tokenize, word_tokenize
from collections import Counter
from PIL import Image
import io

from docx import Document          # DOCX ke liye
from pptx import Presentation      # PPTX ke liye
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4


OCR_AVAILABLE = False
EASYOCR_AVAILABLE = False

try:
    import pytesseract
   
    try:
        import shutil
        found_path = shutil.which('tesseract')
    except Exception:
        found_path = None

    if found_path:
        try:
            pytesseract.pytesseract.tesseract_cmd = found_path
            OCR_AVAILABLE = True
        except Exception:
            OCR_AVAILABLE = False
    else:
        
        common = [r"C:\Program Files\Tesseract-OCR\tesseract.exe", r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe"]
        set_path = False
        for p in common:
            try:
                if os.path.exists(p):
                    pytesseract.pytesseract.tesseract_cmd = p
                    OCR_AVAILABLE = True
                    set_path = True
                    break
            except Exception:
                continue
        if not set_path:
            # no binary found
            pytesseract = None
            OCR_AVAILABLE = False
except ImportError:
    pytesseract = None
    OCR_AVAILABLE = False


try:
    nltk.data.find('tokenizers/punkt_tab')
except LookupError:
    nltk.download('punkt_tab', quiet=True)

try:
    nltk.data.find('corpora/stopwords')
except LookupError:
    nltk.download('stopwords', quiet=True)

STOP_WORDS = set(stopwords.words('english'))

# ---------- TEXT EXTRACTION PDF ----------

def extract_text_from_pdf(pdf_path):
    """Extract text from PDF - handles both typed and handwritten text, including image-based PDFs"""
    import re
    
    doc = fitz.open(pdf_path)
    full_text = ""
    total_pages = len(doc)
    
    for page_num, page in enumerate(doc):
        txt = page.get_text()
        
        # Clean up page numbers and common artifacts
        txt = re.sub(r'^\s*-?\s*\d+\s*-?\s*$', '', txt, flags=re.MULTILINE)  
        txt = re.sub(r'\n\s*\d+\s*$', '\n', txt) 
        
        # If text extraction is minimal, try alternative methods
        if not txt.strip() or len(txt.strip()) < 50:
            # Try raw dict extraction
            try:
                raw_text = page.get_text("rawdict")
                if isinstance(raw_text, dict) and 'blocks' in raw_text:
                    extracted = "\n".join([block.get('text', '') for block in raw_text['blocks'] if block.get('text', '')])
                    if extracted.strip() and len(extracted.strip()) > 20:
                        txt = extracted
            except Exception:
                pass
            
            # If still minimal text, use OCR on image
            if not txt.strip() or len(txt.strip()) < 50:
                try:
                    # Try EasyOCR first (more reliable)
                    try:
                        import easyocr
                        reader = easyocr.Reader(['en'], gpu=False)
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                        image_data = pix.tobytes("ppm")
                        image = Image.open(io.BytesIO(image_data))
                        results = reader.readtext(image, detail=0)
                        if isinstance(results, list):
                            ocr_text = "\n".join([r for r in results if isinstance(r, str) and r.strip()])
                        else:
                            ocr_text = "\n".join(str(r) for r in results if r)
                        if ocr_text.strip():
                            txt = ocr_text
                    except Exception:
                        # Fallback to Tesseract if EasyOCR not available
                        if OCR_AVAILABLE:
                            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                            image_data = pix.tobytes("ppm")
                            image = Image.open(io.BytesIO(image_data))
                            ocr_text = pytesseract.image_to_string(image, lang='eng')
                            if ocr_text.strip():
                                txt = ocr_text
                except Exception:
                    pass
        
        if txt.strip():
            full_text += txt + "\n"

    doc.close()
    
    if not full_text.strip():
        return "Unable to extract text from this PDF. The document may be image-only without readable text."
    
    return full_text

def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    full_text = []
    for para in doc.paragraphs:
        if para.text:
            full_text.append(para.text)
    return "\n".join(full_text)

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n".join(texts)

# ---------- SUMMARIZE (approx 35%) ----------

def summarize_text(text, ratio=0.35, min_sentences=3):
    """Summarize text while prioritizing sentences with dates, numbers, and important keywords"""
    import re
    
    text = text.strip()
    if not text:
        return "", []
    
    sentences = sent_tokenize(text)
    total = len(sentences)
    
    if total == 0:
        return "", []
    
    if total <= 3:
        return text, sentences

    target = max(int(total * ratio), min_sentences)
    if target >= total:
        target = total - 1 if total > 1 else 1

    words = word_tokenize(text.lower())
    words = [w for w in words if w.isalpha() and w not in STOP_WORDS]
    freq = Counter(words)
    
    if not freq:
        return " ".join(sentences[:target]), sentences[:target]

    sentence_scores = {}
    for sent in sentences:
        sent_lower = sent.lower()
        score = 0
        
        for word in freq:
            if word in sent_lower:
                score += freq[word]
        
        if re.search(r'\b\d{1,2}[-/]\d{1,2}[-/]\d{4}\b|\b\d{4}\b', sent):
            score += 50
        
        if re.search(r'\d+', sent):
            score += 20
        
        if re.search(r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)+\b', sent):
            score += 15
        
        sentence_scores[sent] = score

    top_sentences = sorted(
        sentence_scores, key=sentence_scores.get, reverse=True
    )[:target]

    summary_sentences = [s for s in sentences if s in top_sentences]
    summary = " ".join(summary_sentences)
    return summary, summary_sentences

# ---------- DETECT HEADINGS AND SUMMARIZE BY SECTION ----------

def extract_sections_with_summaries(text, ratio=0.35):
    """Extract headings from text and summarize content under each heading"""
    import re
    
    text = text.strip()
    if not text:
        return []
    
    sections = []
    
    heading_patterns = [
        r'^([A-Z][A-Z\s]{3,}?)$',  
        r'^\d+[\.\)]\s+(.+)$',      
        r'^#{1,6}\s+(.+)$',         
        r'^-{3,}$',                 
    ]
    
    lines = text.split('\n')
    current_section = None
    current_content = []
    
    for i, line in enumerate(lines):
        line = line.strip()
        
        # Skip empty lines and pure separators
        if not line or line.startswith('---'):
            continue
        
        # Check if line is a heading
        is_heading = False
        heading_text = None
        
        for pattern in heading_patterns[:-1]:  # Skip separator pattern
            match = re.match(pattern, line)
            if match:
                is_heading = True
                heading_text = match.group(1) if match.groups() else line
                break
        
        if is_heading and heading_text:
            # Save previous section if it exists
            if current_section and current_content:
                content_text = " ".join(current_content).strip()
                if content_text:
                    summary, _ = summarize_text(content_text, ratio=ratio)
                    sections.append({
                        'heading': current_section,
                        'content': content_text,
                        'summary': summary
                    })
            
            # Start new section
            current_section = heading_text
            current_content = []
        else:
            # Add to current section
            if current_section is not None:
                current_content.append(line)
    
    # Don't forget last section
    if current_section and current_content:
        content_text = " ".join(current_content).strip()
        if content_text:
            summary, _ = summarize_text(content_text, ratio=ratio)
            sections.append({
                'heading': current_section,
                'content': content_text,
                'summary': summary
            })
    
    return sections
# [web:32][web:6]

def get_keywords(text, top_n=10):
    """Extract important keywords and detect entities like dates, places, names"""
    import re
    
    words = word_tokenize(text.lower())
    words = [w for w in words if w.isalpha() and w not in STOP_WORDS]
    freq = Counter(words)
    most_common = [w for w, c in freq.most_common(top_n)]
    
    # Detect dates (DD/MM/YYYY, DD-MM-YYYY, YYYY, Month Year)
    date_pattern = r'\b\d{1,2}[-/]\d{1,2}[-/]\d{4}\b|\b\d{4}\b|\b(?:January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{4}\b'
    dates = re.findall(date_pattern, text, re.IGNORECASE)
    
    # Detect places (proper nouns - capitalized words)
    place_pattern = r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b'
    places = list(set(re.findall(place_pattern, text)))[:3]
    
    # Combine all important terms
    important_terms = most_common + dates + places
    return list(set(important_terms))[:15]  # Return top 15 unique terms

# ---------- HIGHLIGHT IN PDF (RED) ----------

def extract_important_terms(text):
    """Extract important terms: acronyms, multi-capital words, names, places, dates, numbers"""
    import re
    important_terms = set()
    
    # 1. Acronyms (2+ capital letters): AI, ML, AWS, USA, API, Q1, Q2, etc
    acronyms = re.findall(r'\b[A-Z]{2,}\b', text)
    important_terms.update(acronyms)
    
    # 2. Words with 2 or more capital letters (SQLServer, MySQL, JavaScript, etc)
    multi_cap_words = re.findall(r'\b[A-Z]\w*(?:[A-Z]\w*)+\b', text)
    important_terms.update(multi_cap_words)
    
    # 3. Proper nouns/Names (Capitalized words at start of sentence or common patterns)
    # Match capitalized words that appear frequently or are person/place names
    proper_nouns = re.findall(r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b', text)
    # Filter for likely names (longer than 1 word or appears multiple times)
    proper_noun_freq = Counter(proper_nouns)
    top_proper = [word for word, count in proper_noun_freq.most_common(10) if count >= 1]
    important_terms.update(top_proper)
    
    # 4. Dates (DD/MM/YYYY, YYYY, Month Year)
    dates = re.findall(
        r'\b\d{1,2}[-/]\d{1,2}[-/]\d{4}\b|\b\d{4}\b|\b(?:January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{4}\b',
        text,
        re.IGNORECASE
    )
    important_terms.update(dates)
    
    # 5. Large numbers and percentages (4+ digits or percentages)
    important_numbers = re.findall(
        r'\b\d{4,}(?:\.\d+)?(?:\s*(?:million|billion|thousand|k|m|b))?\b|\b\d+(?:\.\d+)?%\b',
        text,
        re.IGNORECASE
    )
    important_terms.update(important_numbers)
    
    # 6. Top technical/important keywords (selective list)
    top_tech_terms = [
        'python', 'django', 'javascript', 'database', 'api', 'machine', 'learning',
        'ai', 'ml', 'security', 'cloud', 'aws', 'google', 'deployment', 'server'
    ]
    for term in top_tech_terms:
        if re.search(r'\b' + term + r'\b', text, re.IGNORECASE):
            important_terms.add(term.lower())
    
    return list(important_terms)



def highlight_sentences_with_keywords(input_pdf, output_pdf, summary_sentences, keywords):
    """Highlight important terms in PDF using yellow annotations"""
    import fitz
    import re
    
    # Extract important terms
    try:
        doc = fitz.open(input_pdf)
        full_text = ""
        for page in doc:
            full_text += page.get_text()
        doc.close()
    except Exception:
        full_text = ""
    
    important_terms = extract_important_terms(full_text)
    if not important_terms:
        import shutil
        shutil.copy(input_pdf, output_pdf)
        return

    important_set = set(str(term).lower() for term in important_terms if term)
    
    doc = fitz.open(input_pdf)
    
    for page_num, page in enumerate(doc):
        words = page.get_text("words")
        
        for word_info in words:
            x0, y0, x1, y1, word_text = word_info[0], word_info[1], word_info[2], word_info[3], word_info[4]
            
            should_highlight = False
            
            # Check if word is in important set
            if word_text.lower() in important_set:
                should_highlight = True
            # Check for acronyms (all caps, 2+)
            elif re.match(r'^[A-Z]{2,}$', word_text):
                should_highlight = True
            # Check for multi-capital words (SQLServer, MySQL, etc)
            elif re.match(r'^[A-Z]\w*(?:[A-Z]\w*)+$', word_text):
                should_highlight = True
            # Check for large numbers or percentages
            elif re.match(r'^\d{4,}$|^\d+%$', word_text):
                should_highlight = True
            # Check for proper nouns (starts with capital + lowercase)
            elif re.match(r'^[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*$', word_text) and len(word_text) > 3:
                should_highlight = True
            
            if should_highlight:
                try:
                    rect = fitz.Rect(x0 - 1, y0 - 1, x1 + 1, y1 + 1)
                    highlight = page.add_highlight_annot(rect)
                    highlight.set_colors({"stroke": [1, 1, 0]})  # Yellow
                    highlight.update()
                except Exception:
                    pass
    
    doc.save(output_pdf)
    doc.close()


def highlight_pdf_by_rendering(input_pdf, output_pdf, keywords):
    """Highlight important terms in PDF using direct annotation on PDF pages"""
    import fitz
    import re
    
    # Extract important terms
    try:
        doc = fitz.open(input_pdf)
        full_text = ""
        total_pages = len(doc)
        
        # For large PDFs, sample pages to extract text faster
        page_sample = min(10, total_pages)
        step = max(1, total_pages // page_sample)
        
        for idx in range(0, total_pages, step):
            try:
                p = doc[idx]
                full_text += p.get_text() + "\n"
            except Exception:
                pass
                
        doc.close()
    except Exception:
        full_text = ""

    important_terms = extract_important_terms(full_text)
    if not important_terms:
        # nothing to highlight: copy file
        import shutil
        shutil.copy(input_pdf, output_pdf)
        return

    important_set = set(str(t).lower() for t in important_terms if t)

    # Now open and annotate the PDF
    doc = fitz.open(input_pdf)
    total_pages = len(doc)
    
    for page_no in range(total_pages):
        try:
            page = doc[page_no]
            words = page.get_text("words")
            
            if words:
                for w in words:
                    x0, y0, x1, y1, wtext = w[0], w[1], w[2], w[3], w[4]
                    lw = wtext.strip()
                    if not lw:
                        continue

                    lw_clean = lw.lower()
                    should_highlight = False
                    
                    # Check if word is in important set
                    if lw_clean in important_set:
                        should_highlight = True
                    # Check for acronyms (all caps, 2+)
                    elif re.match(r'^[A-Z]{2,}$', lw):
                        should_highlight = True
                    # Check for multi-capital words (SQLServer, MySQL, etc)
                    elif re.match(r'^[A-Z]\w*(?:[A-Z]\w*)+$', lw):
                        should_highlight = True
                    # Check for large numbers or percentages
                    elif re.match(r'^\d{4,}$|^\d+%$', lw):
                        should_highlight = True
                    # Check for proper nouns (starts with capital + lowercase)
                    elif re.match(r'^[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*$', lw) and len(lw) > 3:
                        should_highlight = True
                    # Check for 2+ uppercase letters in mixed case (camelCase with capitals)
                    elif sum(1 for c in lw if c.isupper()) >= 2:
                        should_highlight = True

                    if should_highlight:
                        try:
                            # Create highlight rectangle with yellow background
                            rect = fitz.Rect(x0 - 1, y0 - 1, x1 + 1, y1 + 1)
                            highlight = page.add_highlight_annot(rect)
                            highlight.set_colors({"stroke": [1, 1, 0]})  # Yellow
                            highlight.update()
                        except Exception:
                            pass
        except Exception:
            pass
    
    doc.save(output_pdf)
    doc.close()
# [web:1][web:5][web:14][web:43]

# ---------- CREATE PROFESSIONAL SUMMARY PDF (ChatGPT STYLE) ----------

def create_professional_pdf_summary(summary_text, output_pdf, keywords=[], sections=None):
    """Create a clean, professional PDF summary with ChatGPT-like formatting and highlighted keywords"""
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
    from reportlab.lib.enums import TA_JUSTIFY, TA_LEFT, TA_CENTER
    from reportlab.lib import colors
    import re
    
    # Create document
    doc = SimpleDocTemplate(
        output_pdf,
        pagesize=A4,
        rightMargin=0.8*inch,
        leftMargin=0.8*inch,
        topMargin=1*inch,
        bottomMargin=0.8*inch
    )
    
    # Get styles
    styles = getSampleStyleSheet()
    
    # Custom styles for professional look
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=20,
        textColor='#1e40af',
        spaceAfter=6,
        alignment=TA_LEFT,
        fontName='Helvetica-Bold',
        leading=24
    )
    
    subtitle_style = ParagraphStyle(
        'Subtitle',
        parent=styles['Normal'],
        fontSize=10,
        textColor='#64748b',
        spaceAfter=16,
        alignment=TA_LEFT,
        fontName='Helvetica-Oblique'
    )
    
    summary_style = ParagraphStyle(
        'Summary',
        parent=styles['BodyText'],
        fontSize=11,
        leading=18,
        textColor='#1f2937',
        alignment=TA_JUSTIFY,
        spaceAfter=12,
        fontName='Helvetica'
    )
    
    # GREEN heading style for section headings
    section_heading_style = ParagraphStyle(
        'SectionHeading',
        parent=styles['Heading2'],
        fontSize=13,
        textColor='#16a34a',  # GREEN color
        spaceAfter=10,
        spaceBefore=12,
        fontName='Helvetica-Bold',
        leading=16
    )
    
    heading_style = ParagraphStyle(
        'PageHeading',
        parent=styles['Heading2'],
        fontSize=13,
        textColor='#0f172a',
        spaceAfter=10,
        spaceBefore=12,
        fontName='Helvetica-Bold',
        leading=16
    )
    
    def highlight_keywords_in_text(text, keywords):
        """Add highlights for keywords, dates, numbers in text"""
        # Escape HTML characters
        processed = text.replace('&', '&amp;')
        processed = processed.replace('<', '&lt;')
        processed = processed.replace('>', '&gt;')
        
        # Highlight keywords - CYAN/BLUE with background
        if keywords:
            keywords_clean = [kw.lower() for kw in keywords if kw]
            for kw in keywords_clean:
                # Use word boundaries to match whole words
                pattern = r'\b' + re.escape(kw) + r'\b'
                processed = re.sub(
                    pattern,
                    lambda m: f'<b><font color="#0284c7" bgcolor="#e0f2fe">{m.group(0)}</font></b>',
                    processed,
                    flags=re.IGNORECASE
                )
        
        # Highlight dates - RED and bold
        processed = re.sub(
            r'\b(\d{1,2}[-/]\d{1,2}[-/]\d{4}|\d{4})\b',
            r'<b><font color="#dc2626">\1</font></b>',
            processed
        )
        
        # Highlight numbers/percentages - ORANGE
        processed = re.sub(
            r'\b(\d+(?:\.\d+)?%?)\b',
            r'<b><font color="#ea580c">\1</font></b>',
            processed
        )
        
        return processed
    
    # Build story
    story = []
    
    # Add header
    story.append(Paragraph("📄 Document Summary", title_style))
    # story.append(Paragraph("AI-Generated Summary - Key Information Extracted", subtitle_style))
    story.append(Spacer(1, 0.1*inch))
    
    # If sections are provided, show section-based summary
    if sections and len(sections) > 0:
        story.append(Paragraph("📋 Section-wise Summary", heading_style))
        story.append(Spacer(1, 0.08*inch))
        
        for section in sections:
            # Add GREEN heading
            heading_text = section.get('heading', 'Untitled').upper()
            story.append(Paragraph(f"✓ <b><font color=\"#16a34a\">{heading_text}</font></b>", section_heading_style))
            
            # Add summary for this section with highlights
            section_summary = section.get('summary', '').strip()
            if section_summary:
                processed = highlight_keywords_in_text(section_summary, keywords)
                try:
                    story.append(Paragraph(processed, summary_style))
                except Exception:
                    story.append(Paragraph(section_summary, summary_style))
            
            story.append(Spacer(1, 0.1*inch))
    else:
        # Fallback to original summary if no sections
        story.append(Paragraph("📝 Summary", heading_style))
        
        processed_summary = highlight_keywords_in_text(summary_text, keywords)
        
        sentences = processed_summary.split('.')
        current_para = ""
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
            
            current_para += sentence + ". "
            
            if current_para.count('.') >= 2:
                try:
                    story.append(Paragraph(current_para.strip(), summary_style))
                    current_para = ""
                except Exception:
                    current_para = ""
        
        if current_para.strip():
            try:
                story.append(Paragraph(current_para.strip(), summary_style))
            except Exception:
                pass
    
    story.append(Spacer(1, 0.15*inch))
    
    # Add keywords section
    if keywords:
        story.append(Paragraph("Key Points", heading_style))
        
        keywords_list = [kw for kw in keywords if kw][:12]
        
        if keywords_list:
            keywords_text = " • ".join([f"<b><font color=\"#0284c7\">{kw}</font></b>" for kw in keywords_list])
            try:
                story.append(Paragraph(keywords_text, summary_style))
            except Exception:
                pass
            story.append(Spacer(1, 0.1*inch))
    
    # Add footer
    story.append(Spacer(1, 0.2*inch))
    footer_style = ParagraphStyle(
        'Footer',
        parent=styles['Normal'],
        fontSize=8,
        textColor='#94a3b8',
        alignment=TA_CENTER,
        fontName='Helvetica'
    )
    story.append(Paragraph("Generated by dipendra minor project | 📄 AI-Powered Analysis", footer_style))
    
    # Build PDF
    try:
        doc.build(story)
    except Exception as e:
        # Fallback: create simple text PDF if fancy formatting fails
        doc.build([
            Paragraph("📄 Document Summary", title_style),
            Paragraph("AI-Generated Summary", subtitle_style),
            Paragraph(re.sub(r'<[^>]+>', '', summary_text), summary_style)
        ])
# [web:6][web:23]

# ---------- PIPELINES ----------

def process_pdf(input_path, output_path):
    """Process PDF: extract text, summarize, and create SUMMARY PDF with highlights"""
    import os
    
    text = extract_text_from_pdf(input_path)
    
    # Extract sections with summaries
    sections = extract_sections_with_summaries(text, ratio=0.25)  # 25% for more condensed summary
    
    # If sections found, use them; otherwise fall back to overall summary
    if sections and len(sections) > 0:
        summary = " ".join([s.get('summary', '') for s in sections if s.get('summary')])
    else:
        summary, _ = summarize_text(text, ratio=0.25)  # 25% summary
    
    keywords = get_keywords(text, top_n=10)
    
    # Create professional summary PDF (SHORT version)
    try:
        create_professional_pdf_summary(summary, output_path, keywords, sections=sections)
    except Exception as e:
        # Fallback: if summary PDF creation fails, create highlighted version
        try:
            highlight_pdf_by_rendering(input_path, output_path, keywords)
        except Exception:
            try:
                highlight_sentences_with_keywords(input_path, output_path, [], keywords)
            except Exception:
                # Last resort: copy original
                import shutil
                shutil.copy(input_path, output_path)

    return summary, keywords

def process_docx(input_path, output_path):
    text = extract_text_from_docx(input_path)
    
    # Extract sections with summaries
    sections = extract_sections_with_summaries(text, ratio=0.35)
    
    if sections and len(sections) > 0:
        summary = " ".join([s.get('summary', '') for s in sections if s.get('summary')])
    else:
        summary, _ = summarize_text(text, ratio=0.35)
    
    keywords = get_keywords(summary, top_n=10)
    
    create_professional_pdf_summary(summary, output_path, keywords, sections=sections)
    return summary, keywords

def process_pptx(input_path, output_path):
    text = extract_text_from_pptx(input_path)
    
    # Extract sections with summaries
    sections = extract_sections_with_summaries(text, ratio=0.35)
    
    if sections and len(sections) > 0:
        summary = " ".join([s.get('summary', '') for s in sections if s.get('summary')])
    else:
        summary, _ = summarize_text(text, ratio=0.35)
    
    keywords = get_keywords(summary, top_n=10)

    create_professional_pdf_summary(summary, output_path, keywords, sections=sections)
    return summary, keywords

def process_file_by_extension(input_path, output_path):
    ext = os.path.splitext(input_path)[1].lower()
    if ext == ".pdf":
        return process_pdf(input_path, output_path)
    elif ext == ".docx":
        return process_docx(input_path, output_path)
    elif ext in [".pptx", ".ppt"]:
        return process_pptx(input_path, output_path)
    else:
        raise ValueError("Unsupported file type")

