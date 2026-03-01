import os
import io
from pathlib import Path
from typing import List, Tuple

from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.pdfgen import canvas
from reportlab.lib.units import mm

try:
    from docx import Document
except Exception:
    Document = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from pdfminer.high_level import extract_text as pdf_extract_text
except Exception:
    pdf_extract_text = None

from sklearn.feature_extraction.text import TfidfVectorizer


def extract_text_from_docx(path: str) -> str:
    if Document is None:
        raise RuntimeError('python-docx not installed')
    doc = Document(path)
    return '\n'.join(p.text for p in doc.paragraphs)


def extract_text_from_pptx(path: str) -> str:
    if Presentation is None:
        raise RuntimeError('python-pptx not installed')
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                texts.append(shape.text)
    return '\n'.join(texts)


def extract_text_from_pdf(path: str) -> str:
    if pdf_extract_text is None:
        raise RuntimeError('pdfminer.six not installed')
    return pdf_extract_text(path)


def split_sentences(text: str) -> List[str]:
    # Very small sentence splitter by punctuation
    import re
    sentences = re.split(r'(?<=[.!?])\s+', text.replace('\r', '\n'))
    return [s.strip() for s in sentences if s.strip()]


def score_sentences(sentences: List[str], top_k_words: int = 20) -> Tuple[List[float], List[str]]:
    if not sentences:
        return [], []
    vec = TfidfVectorizer(stop_words='english')
    tfidf = vec.fit_transform(sentences)
    # Sentence score: sum of tf-idf weights in the sentence
    scores = tfidf.sum(axis=1).A1.tolist()
    # Top words across the full set of sentences
    tfidf_full = TfidfVectorizer(stop_words='english', max_features=top_k_words)
    tfidf_full.fit([' '.join(sentences)])
    top_words = tfidf_full.get_feature_names_out().tolist()
    return scores, top_words


def create_highlighted_pdf(text: str, sentence_scores: List[float], top_words: List[str], out_path: str):
    # Render text to PDF and highlight sentences that are above median score and words that match top_words
    sentences = split_sentences(text)
    if not sentences:
        raise ValueError('No text to write')

    median = sorted(sentence_scores)[len(sentence_scores)//2] if sentence_scores else 0

    c = canvas.Canvas(out_path, pagesize=A4)
    width, height = A4
    left_margin = 20 * mm
    right_margin = 20 * mm
    y = height - 20 * mm
    max_width = width - left_margin - right_margin
    line_height = 10 * mm / 3

    c.setFont('Helvetica', 10)

    for i, sent in enumerate(sentences):
        words = sent.split()
        # prepare line-wrapping
        line = ''
        for w in words:
            test = (line + ' ' + w).strip()
            tw = c.stringWidth(test, 'Helvetica', 10)
            if tw > max_width:
                # draw current line
                if sentence_scores and sentence_scores[i] >= median:
                    # highlight line background pink
                    c.setFillColor(colors.Color(1, 0.85, 0.9))
                    c.rect(left_margin, y - line_height*0.25, max_width, line_height, stroke=0, fill=1)
                    c.setFillColor(colors.black)
                # draw text with highlighted words
                draw_text_with_word_highlights(c, left_margin + 2, y - line_height*0.1, line, top_words)
                y -= line_height
                line = w
                if y < 20*mm:
                    c.showPage()
                    c.setFont('Helvetica', 10)
                    y = height - 20 * mm
            else:
                line = test

        # draw remaining line
        if line:
            if sentence_scores and sentence_scores[i] >= median:
                c.setFillColor(colors.Color(1, 0.85, 0.9))
                c.rect(left_margin, y - line_height*0.25, max_width, line_height, stroke=0, fill=1)
                c.setFillColor(colors.black)
            draw_text_with_word_highlights(c, left_margin + 2, y - line_height*0.1, line, top_words)
            y -= line_height
            if y < 20*mm:
                c.showPage()
                c.setFont('Helvetica', 10)
                y = height - 20 * mm

    c.save()


def draw_text_with_word_highlights(c: canvas.Canvas, x: float, y: float, text: str, top_words: List[str]):
    # Draw text word by word, coloring top words pink
    words = text.split()
    cur_x = x
    for w in words:
        spacer = c.stringWidth(' ', 'Helvetica', 10)
        w_clean = ''.join(ch for ch in w if ch.isalnum()).lower()
        w_width = c.stringWidth(w, 'Helvetica', 10)
        if w_clean in top_words:
            c.setFillColor(colors.HexColor('#ff66b2'))  # pink for important words
        else:
            c.setFillColor(colors.black)
        c.drawString(cur_x, y, w)
        cur_x += w_width + spacer
    c.setFillColor(colors.black)


def summarize_and_highlight(input_path: str, output_path: str = None):
    p = Path(input_path)
    if not p.exists():
        raise FileNotFoundError(input_path)
    suffix = p.suffix.lower()
    if suffix == '.pdf':
        text = extract_text_from_pdf(str(p))
    elif suffix in ('.docx',):
        text = extract_text_from_docx(str(p))
    elif suffix in ('.pptx',):
        text = extract_text_from_pptx(str(p))
    else:
        # fallback: read as plain text
        text = p.read_text(encoding='utf-8', errors='ignore')

    sentences = split_sentences(text)
    scores, top_words = score_sentences(sentences, top_k_words=30)

    if output_path is None:
        out_dir = Path('media/processed')
        out_dir.mkdir(parents=True, exist_ok=True)
        output_path = str(out_dir / (p.stem + '_summary_highlighted.pdf'))

    create_highlighted_pdf(text, scores, top_words, output_path)
    return output_path


if __name__ == '__main__':
    import sys
    if len(sys.argv) < 2:
        print('Usage: python highlight_summarizer.py <input-file> [output-file]')
        sys.exit(1)
    inp = sys.argv[1]
    out = sys.argv[2] if len(sys.argv) > 2 else None
    print('Processing', inp)
    outp = summarize_and_highlight(inp, out)
    print('Saved to', outp)
